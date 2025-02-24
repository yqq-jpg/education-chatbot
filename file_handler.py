import os
from typing import Dict, Any
import pandas as pd
from docx import Document
import PyPDF2
import chardet
from pptx import Presentation

class FileHandler:
    def __init__(self):
        self.supported_extensions = {
            'txt': self._process_txt,
            'docx': self._process_docx,
            'pdf': self._process_pdf,
            'xlsx': self._process_excel,
            'xls': self._process_excel,
            'csv': self._process_csv,
            'pptx': self._process_pptx
        }

    def process_file(self, file_path: str) -> Dict[str, Any]:
        try:
            # 基本检查
            if not os.path.exists(file_path):
                return {'status': 'error', 'message': 'File not found'}
            
            if os.path.getsize(file_path) == 0:
                return {'status': 'error', 'message': 'File is empty'}

            original_name = os.path.basename(file_path)
            file_extension = ''

            # 获取扩展名
            if '.' in original_name:
                file_extension = original_name.rsplit('.', 1)[1].lower()
            else:
                path_parts = file_path.lower().split('\\')
                for part in path_parts:
                    if part in self.supported_extensions:
                        file_extension = part
                        break

            # 如果仍然没有扩展名，根据路径最后一部分判断
            if not file_extension:
                last_part = path_parts[-1]
                if last_part in self.supported_extensions:
                    file_extension = last_part

            print(f"Processing file: {file_path}, detected extension: {file_extension}")

            if not file_extension or file_extension not in self.supported_extensions:
                return {
                    'status': 'error', 
                    'message': f'Unsupported or missing file type: {file_extension}'
                }

            # 处理文件
            processor = self.supported_extensions[file_extension]
            result = processor(file_path)
            
            return {
                'status': 'success',
                'content': result['content'],
                'metadata': result['metadata']
            }

        except Exception as e:
            print(f"Detailed error processing file: {str(e)}")
            return {'status': 'error', 'message': str(e)}

    def _process_txt(self, file_path: str) -> Dict[str, Any]:
        try:
            # 首先尝试使用 chardet 检测编码
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                result = chardet.detect(raw_data)
                encoding = result['encoding'] if result['encoding'] else 'utf-8'

            with open(file_path, 'r', encoding=encoding) as file:
                content = file.read()

            return {
                'content': content,
                'metadata': {
                    'file_size': os.path.getsize(file_path),
                    'encoding': encoding,
                    'line_count': len(content.splitlines())
                }
            }
        except UnicodeDecodeError:
            # 如果检测到的编码不正确，尝试其他常用编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'ascii', 'iso-8859-1']
            for enc in encodings:
                try:
                    with open(file_path, 'r', encoding=enc) as file:
                        content = file.read()
                        return {
                            'content': content,
                            'metadata': {
                                'file_size': os.path.getsize(file_path),
                                'encoding': enc,
                                'line_count': len(content.splitlines())
                            }
                        }
                except UnicodeDecodeError:
                    continue
            raise ValueError("Unable to decode file with any supported encoding")

    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        try:
            doc = Document(file_path)
            content = []

            # 处理段落
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)

            # 处理表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        content.append(' | '.join(row_text))

            return {
                'content': '\n'.join(content),
                'metadata': {
                    'file_size': os.path.getsize(file_path),
                    'paragraph_count': len(doc.paragraphs),
                    'table_count': len(doc.tables)
                }
            }
        except Exception as e:
            raise ValueError(f"Error processing DOCX file: {str(e)}")

    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                content = []

                for page in reader.pages:
                    text = page.extract_text()
                    if text.strip():
                        content.append(text)

                return {
                    'content': '\n\n'.join(content),
                    'metadata': {
                        'file_size': os.path.getsize(file_path),
                        'page_count': len(reader.pages)
                    }
                }
        except Exception as e:
            raise ValueError(f"Error processing PDF file: {str(e)}")

    def _process_excel(self, file_path: str) -> Dict[str, Any]:
        try:
            print(f"Processing Excel file: {file_path}")
            df = pd.read_excel(file_path, engine='openpyxl')
            try:
                content = df.to_string()
                if isinstance(content, bytes):
                    content = content.decode('utf-8')
            except:
                content = df.to_string(encoding='utf-8')

            return {
                'content': content,
                'metadata': {
                    'file_size': os.path.getsize(file_path),
                    'row_count': len(df),
                    'column_count': len(df.columns)
                }
            }
        except Exception as e:
            print(f"Excel processing error: {str(e)}")
            raise ValueError(f"Error processing Excel file: {str(e)}")

    def _process_csv(self, file_path: str) -> Dict[str, Any]:
        try:
            # 首先尝试使用 chardet 检测编码
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                result = chardet.detect(raw_data)
                encoding = result['encoding'] if result['encoding'] else 'utf-8'

            df = pd.read_csv(file_path, encoding=encoding)
            
            return {
                'content': df.to_string(),
                'metadata': {
                    'file_size': os.path.getsize(file_path),
                    'row_count': len(df),
                    'column_count': len(df.columns),
                    'encoding': encoding
                }
            }
        except Exception as e:
            encodings = ['utf-8', 'gbk', 'gb2312', 'iso-8859-1']
            for enc in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=enc)
                    return {
                        'content': df.to_string(),
                        'metadata': {
                            'file_size': os.path.getsize(file_path),
                            'row_count': len(df),
                            'column_count': len(df.columns),
                            'encoding': enc
                        }
                    }
                except Exception:
                    continue
            raise ValueError(f"Error processing CSV file: {str(e)}")

    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        try:
            prs = Presentation(file_path)
            content = []

            for slide in prs.slides:
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                if slide_content:
                    content.append('\n'.join(slide_content))

            return {
                'content': '\n\n'.join(content),
                'metadata': {
                    'file_size': os.path.getsize(file_path),
                    'slide_count': len(prs.slides)
                }
            }
        except Exception as e:
            raise ValueError(f"Error processing PPTX file: {str(e)}")